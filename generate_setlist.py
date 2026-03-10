import os
import argparse
import configparser
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io

def get_config():
    config = configparser.ConfigParser()
    config_path = os.path.expanduser('~/.config/jtgs-ppt/jtgs.conf')
    
    # Defaults if config file is missing
    defaults = {
        'song_folder': './songs',
        'intro_slide': './intro.pptx',
        'transition_slide': './transition.pptx'
    }
    
    if os.path.exists(config_path):
        config.read(config_path)
        return {
            'song_folder': config.get('Paths', 'song_folder', fallback=defaults['song_folder']),
            'intro_slide': config.get('Paths', 'intro_slide', fallback=defaults['intro_slide']),
            'transition_slide': config.get('Paths', 'transition_slide', fallback=defaults['transition_slide'])
        }
    return defaults

def append_slides(source_prs, target_prs):
    for slide in source_prs.slides:
        # Use a blank layout (index 6) from the target's master
        slide_layout = target_prs.slide_layouts[6]
        new_slide = target_prs.slides.add_slide(slide_layout)
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # For pictures, we must copy the image blob and add it properly
                # to the target slide to handle relationships and media parts.
                image_stream = io.BytesIO(shape.image.blob)
                new_slide.shapes.add_picture(
                    image_stream, 
                    shape.left, 
                    shape.top, 
                    shape.width, 
                    shape.height
                )
            else:
                # For other shapes, attempt to copy the XML element directly.
                # Note: This may still have issues with shapes that have relationships.
                new_slide.shapes._spTree.append(shape.element)

def main():
    cfg = get_config()
    today = datetime.now().strftime("%Y_%m_%d")
    default_out = f"JTGS_{today}.pptx"

    parser = argparse.ArgumentParser(description="Stitch PPTX songs into a setlist.")
    parser.add_argument('-s', '--songs', nargs='+', required=True, help='Song filenames (no extension)')
    parser.add_argument('-o', '--output', default=default_out, help=f'Output filename (default: {default_out})')
    args = parser.parse_args()

    if not os.path.exists(cfg['intro_slide']):
        print(f"Error: Intro not found at {cfg['intro_slide']}")
        return
    
    master_prs = Presentation(cfg['intro_slide'])
    
    transition_prs = None
    if os.path.exists(cfg['transition_slide']):
        transition_prs = Presentation(cfg['transition_slide'])
    else:
        print(f"Warning: Transition slide missing at {cfg['transition_slide']}")

    for song_name in args.songs:
        song_path = os.path.join(cfg['song_folder'], f"{song_name}.pptx")
        
        if os.path.exists(song_path):
            print(f"Merging: {song_name}")
            append_slides(Presentation(song_path), master_prs)
            if transition_prs:
                append_slides(transition_prs, master_prs)
        else:
            print(f"Skipping: {song_name}.pptx (File not found)")

    master_prs.save(args.output)
    print(f"\nSaved setlist to: {args.output}")

if __name__ == "__main__":
    main()