import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Configuration
INPUT_FOLDER = '/mnt/data/WORK/JTGS/MUSIC/ALL_SONGS'
OUTPUT_FOLDER = '/mnt/data/WORK/JTGS/MUSIC/ALL_SONGS_CONVERTED'

FONT_MAP = {
    "Arial Black": "Noto Sans Black",
    "Default": "Noto Sans"
}

def remove_transitions(slide):
    """Safely removes the transition XML element from its parent."""
    element = slide._element
    # Find the transition element and its parent
    transition = element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    if transition is not None:
        parent = transition.getparent()
        parent.remove(transition)

def apply_font_logic(run):
    current_font = run.font.name
    if current_font == "Arial Black":
        run.font.name = FONT_MAP["Arial Black"]
    else:
        run.font.name = FONT_MAP["Default"]

def process_shapes(shapes):
    """Recursively processes shapes, groups, and tables."""
    for shape in shapes:
        # 1. Handle Text Frames (Standard shapes)
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    apply_font_logic(run)

        # 2. Handle Tables
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            apply_font_logic(run)

        # 3. Handle Grouped Shapes (Recursive)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            process_shapes(shape.shapes)

def run_batch_process():
    if not os.path.exists(OUTPUT_FOLDER):
        os.makedirs(OUTPUT_FOLDER)

    for filename in os.listdir(INPUT_FOLDER):
        if filename.endswith(".pptx") and not filename.startswith("~$"):
            try:
                prs = Presentation(os.path.join(INPUT_FOLDER, filename))

                # Slide Masters & Layouts
                for master in prs.slide_masters:
                    process_shapes(master.shapes)
                    for layout in master.slide_layouts:
                        process_shapes(layout.shapes)

                # Individual Slides
                for slide in prs.slides:
                    remove_transitions(slide)
                    process_shapes(slide.shapes)

                prs.save(os.path.join(OUTPUT_FOLDER, filename))
                print(f"Successfully processed: {filename}")
            except Exception as e:
                print(f"Failed to process {filename}: {e}")

if __name__ == "__main__":
    run_batch_process()
